import os
import time
from datetime import datetime
import traceback # allows error reporting
from multiprocessing import Pool, cpu_count # allows parallelized computing
from tqdm import tqdm # allows visualizing state of a process
from whoosh.fields import Schema, TEXT, NUMERIC # needed to create whoosh index
from whoosh.index import create_in # needed to create whoosh index
import tkinter as tk # needed for GUI-interface to select files
from tkinter import filedialog # needed for GUI-interface to select files
from PyPDF2 import PdfReader # needed to extract pdf data
import pptx # needed to process ppt-files
from markdown2 import markdown_path # needed to process markdown-files
from nbconvert import HTMLExporter # needed to process Jupyter Notebook files

# Initialize variables
unsupported_files = []
error_logs = []
supported_filetypes = [".pdf", ".txt", ".csv", ".py", ".ipynb", ".html", ".r", ".qmd", ".pptx"]

# Function to extract metadata
def extract_metadata(filepath):
    try:
        stat = os.stat(filepath)
        # create dictionary entry for extracted metadata
        metadata = {
            "Filename": os.path.basename(filepath),
            "Path": filepath,
            "Author": "Unknown",
            "CreateDate": datetime.fromtimestamp(stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S"),
            "Pages": 0
        }
        return metadata
    # log errors in log-file
    except Exception as e:
        log_error(filepath, "Metadata extraction failed: " + str(e))
        return None

# Function to log errors
def log_error(filepath, message):
    error_logs.append(f"Error processing {filepath}: {message}")
    unsupported_files.append(filepath)

# Worker function to process a single file
def process_file(filepath):
    global supported_filetypes
    # extract file extension
    file_ext = os.path.splitext(filepath)[1].lower()
    # call metadata extraction function
    metadata = extract_metadata(filepath)
    if not metadata:
        return None

    # log error, if filetype not supported
    if file_ext not in supported_filetypes:
        log_error(filepath, "Unsupported file type")
        return None

    results = [] # object to store index-data
    try:
        # extraction function for PDF files
        if file_ext == ".pdf":
            reader = PdfReader(filepath)
            metadata["Pages"] = len(reader.pages)
            metadata["Author"] = reader.metadata.get("/Author", "Unknown")
            for page_number, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                results.append({
                    "file_name": str(metadata["Filename"]),
                    "path": str(metadata["Path"]),
                    "author": str(metadata["Author"]),
                    "create_date": str(metadata["CreateDate"]),
                    "page": int(page_number),
                    "content": str(text.strip())
                })

        # extraction function for TXT, CSV, PY, HTML files
        elif file_ext in [".txt", ".csv", ".py", ".html"]:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                results.append({
                    "file_name": metadata["Filename"],
                    "path": metadata["Path"],
                    "author": metadata["Author"],
                    "create_date": metadata["CreateDate"],
                    "page": 1,
                    "content": text.strip()
                })

        # extraction function for Jupyter Notebook files
        elif file_ext == ".ipynb":
            with open(filepath, "r", encoding="utf-8") as f:
                notebook_content = f.read()
                exporter = HTMLExporter()
                html_content, _ = exporter.from_notebook_node(notebook_content)
                results.append({
                    "file_name": metadata["Filename"],
                    "path": metadata["Path"],
                    "author": metadata["Author"],
                    "create_date": metadata["CreateDate"],
                    "page": 1,
                    "content": html_content.strip()
                })

        # extraction function for R files
        elif file_ext == ".r":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                results.append({
                    "file_name": metadata["Filename"],
                    "path": metadata["Path"],
                    "author": metadata["Author"],
                    "create_date": metadata["CreateDate"],
                    "page": 1,
                    "content": text.strip()
                })

        # extraction function for PDF files
        elif file_ext == ".qmd":
            html_content = markdown_path(filepath)
            results.append({
                "file_name": metadata["Filename"],
                "path": metadata["Path"],
                "author": metadata["Author"],
                "create_date": metadata["CreateDate"],
                "page": 1,
                "content": html_content.strip()
            })

        # extraction function for PPTX files
        elif file_ext == ".pptx":
            presentation = pptx.Presentation(filepath)
            slide_texts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_texts.append(shape.text)
            results.append({
                "file_name": metadata["Filename"],
                "path": metadata["Path"],
                "author": metadata["Author"],
                "create_date": metadata["CreateDate"],
                "page": 1,
                "content": "\n".join(slide_texts).strip()
            })

        return results

    # Log errors in log-file
    except Exception as e:
        log_error(filepath, traceback.format_exc())
        return None

# Main function to process files in parallel
def process_files_parallel(folders_to_index, index_dir):
    files_to_process = []
    for folder_path in folders_to_index:
        for root, _, files in os.walk(folder_path):
            for file in files:
                files_to_process.append(os.path.join(root, file))

    # defines amount of used processors for parallelizing (max - 2), adjust if needed
    num_workers = max(1, cpu_count() - 2)
    print(f"Processing {len(files_to_process)} files using {num_workers} workers...")

    # Define schema for whoosh index, can be extended, if more metadata fields needed
    schema = Schema(
        file_name=TEXT(stored=True),
        path=TEXT(stored=True),
        author=TEXT(stored=True),
        create_date=TEXT(stored=True),
        page=NUMERIC(stored=True),
        content=TEXT(stored=True)
    )

    # Create Whoosh index
    if not os.path.exists(index_dir):
        os.mkdir(index_dir)
    ix = create_in(index_dir, schema)

    results = []
    # iterate through document database and process parallelized
    with Pool(num_workers) as pool:
        for result in tqdm(pool.imap_unordered(process_file, files_to_process), total=len(files_to_process), desc="Indexing Files", dynamic_ncols=True):
            if result:
                results.extend(result)

    # Write results to Whoosh index
    with ix.writer() as writer:
        for doc in results:
            try:
                # ignore documents without content and create error log entry
                if not doc.get("content"):
                    log_error(doc.get("path", "Unknown"), "Document content is empty. Skipping.")
                    continue
                writer.add_document(**doc)
            except Exception as e:
                log_error(doc.get("path", "Unknown"), str(e))

    # Write error logs, change error file name if needed
    with open(os.path.join(index_dir, "error_log.txt"), "w", encoding="utf-8") as f:
        f.write("\n".join(error_logs))

# Select multiple directories with documents manually
def select_multiple_directories():
    root = tk.Tk()
    root.withdraw()
    folders = []
    while True:
        folder_path = filedialog.askdirectory(title="Select a Folder to Index")
        if not folder_path:
            break
        folders.append(folder_path)
        add_more = tk.messagebox.askyesno("Add More Folders?", "Do you want to add another folder?")
        if not add_more:
            break
    return folders

# Main program
if __name__ == "__main__":
    # get folders to index
    folders_to_index = select_multiple_directories()
    if not folders_to_index:
        print("No folders selected for indexing. Exiting.")
        exit()
    else:
        # get folder to save the index (same path as configured in search function)
        output_dir = filedialog.askdirectory(title="Select a Folder to Save the Index")
        if not output_dir:
            print("No folder selected for saving the index. Exiting.")
        else:
            t = time.time()
            # call processing function
            process_files_parallel(folders_to_index, output_dir)

            print("Indexing complete.")
            print(f"Index saved to: {output_dir}")
            if unsupported_files:
                print("Unsupported files:")
                for file in unsupported_files:
                    print(file)
            print("Error logs saved to error_log.txt in the index directory.")
            print("Processing time:", round(time.time()-t,2), "s")

'''
###
Output by indexing FULL DATABASE of five CAS courses
Running on my private computer (Michel)
###
PS P:\PY\CAS IR Leistungsnachweis> & C:/Python312/python.exe "p:/PY/CAS IR Leistungsnachweis/directory_indexer.py"
Processing 3776 files using 14 workers...
Indexing Files:   9%|████████▌                                                                                  | 356/3776 [00:16<00:56, 60.90it/s]unknown widths : 
[0, IndirectObject(23, 0, 2249369934048)]
unknown widths : 
[0, IndirectObject(25, 0, 2249369934048)]
Indexing Files:  11%|█████████▋                                                                                 | 403/3776 [00:17<00:55, 60.95it/s] impossible to decode XFormObject /Im70
Indexing Files:  53%|███████████████████████████████████████████████▌                                          | 1995/3776 [00:55<01:41, 17.51it/s]Multiple definitions in dictionary at byte 0x22ec6 for key /Author
Multiple definitions in dictionary at byte 0x22ece for key /Title
Indexing Files: 100%|██████████████████████████████████████████████████████████████████████████████████████████| 3776/3776 [02:45<00:00, 22.78it/s]
Indexing complete.
Index saved to: P:/PY/CAS IR Leistungsnachweis/whoosh_index
Unsupported files:
G:/OneDrive - Flex/4.1_CAS Datenanalyse\code.R
...
G:/OneDrive - Flex/4.1_CAS Datenanalyse\datasciencebook_ocred.pdf
Error logs saved to error_log.txt in the index directory.
Processing time: 390.35 s 
'''